"""
DMS – Dokumentenmanagementsystem für Ilija Public Edition v6
=============================================================
Fixes:
- EXIF-Rotation für gestochen scharfes OCR von Handyfotos (ImageOps.exif_transpose)
- Kugelsicherer Pipe-Scanner: findet jede Zeile mit 3+ Trennzeichen, egal ob Sternchen oder Leerzeichen
- Komplette KI-Antwort im Terminal (kein Abschneiden mehr)
- Kurzer Direkt-Prompt bei leerem OCR-Text (kein CoT der Format bricht)
- dms_archiv_baum, dms_pfad_setzen, dms_passwort_entfernen vollständig erhalten
"""

import os
import re
import shutil
import hashlib
import json
from pathlib import Path
from datetime import datetime

# ── Konfiguration ─────────────────────────────────────────────
DMS_BASE_DEFAULT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "dms")

# Erlaubte Dateiendungen (Sicherheitsfilter: keine ausfuehrbaren Dateien)
ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".txt", ".md", ".csv", ".rtf", ".odt", ".ods", ".odp",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp",
    ".heic", ".heif",
    ".eml", ".msg",
    ".xml", ".json", ".html", ".htm",
}

def _get_config() -> dict:
    config_path = os.path.join(DMS_BASE_DEFAULT, "dms_config.json")
    defaults = {
        "archiv_pfad":    os.path.join(DMS_BASE_DEFAULT, "archiv"),
        "import_pfad":    os.path.join(DMS_BASE_DEFAULT, "import"),
        "passwort_hash":  "",
        "passwort_aktiv": False,
    }
    try:
        if os.path.exists(config_path):
            with open(config_path, "r", encoding="utf-8") as f:
                defaults.update(json.load(f))
    except Exception:
        pass
    return defaults

def _save_config(config: dict):
    os.makedirs(DMS_BASE_DEFAULT, exist_ok=True)
    config_path = os.path.join(DMS_BASE_DEFAULT, "dms_config.json")
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)

def _get_archiv_dir() -> str:
    return _get_config()["archiv_pfad"]

def _get_import_dir() -> str:
    return _get_config()["import_pfad"]

def _get_meta_file() -> str:
    return os.path.join(DMS_BASE_DEFAULT, "meta.json")

MAX_TEXT_LEN = 3000

SUPPORTED_EXTS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".xlsm",
    ".txt", ".csv", ".md", ".rtf",
    ".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif",
    ".bmp", ".heic", ".heif",
    ".odt", ".ods", ".odp", ".pptx", ".ppt",
}

# ── Hilfsfunktionen ───────────────────────────────────────────

def _init_dirs():
    cfg = _get_config()
    os.makedirs(cfg["import_pfad"], exist_ok=True)
    os.makedirs(cfg["archiv_pfad"], exist_ok=True)
    os.makedirs(DMS_BASE_DEFAULT, exist_ok=True)
    if not os.path.exists(_get_meta_file()):
        _save_meta({})

def _load_meta() -> dict:
    try:
        with open(_get_meta_file(), "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def _save_meta(meta: dict):
    with open(_get_meta_file(), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

def _berechne_hash(filepath: str) -> str:
    hasher = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
    except Exception:
        return ""
    return hasher.hexdigest()

def _naechste_version(ziel_pfad: str) -> str:
    if not os.path.exists(ziel_pfad):
        return ziel_pfad
    basis, endung = os.path.splitext(ziel_pfad)
    basis = re.sub(r'_v\d+$', '', basis)
    counter = 2
    while True:
        kandidat = f"{basis}_v{counter}{endung}"
        if not os.path.exists(kandidat):
            return kandidat
        counter += 1

def _sanitize(text: str) -> str:
    text = text.strip()
    text = re.sub(r'[\\/*?:"<>|]', '', text)
    text = re.sub(r'\s+', '_', text)
    if text.isupper() and len(text) > 6:
        text = text.capitalize()
    return text[:80]

def _sanitize_filename(text: str, endung: str) -> str:
    text = text.strip()
    if text.lower().endswith(endung.lower()):
        text = text[:-len(endung)]
    text = re.sub(r'[\\/*?:"<>|]', '', text)
    text = re.sub(r'\s+', '_', text)
    text = text.strip('._-')
    verdaechtig = {'neuer_dateiname', 'dateiname', 'filename', 'neuer', 'name', 'dokument', 'unbekannt', 'unknown'}
    if not text or text.lower() in verdaechtig or len(text) < 3:
        text = f"Dokument_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    return text[:80] + endung

def _is_kryptisch(filename: str) -> bool:
    stem = Path(filename).stem.lower()
    if len(stem) < 5: return True
    if re.match(r'^[0-9a-f]{8,}$', stem): return True
    if re.match(r'^[0-9_\-]+$', stem): return True
    if re.match(r'^(img|image|photo|foto|scan|doc|file|dokument|unnamed|untitled|new|neu)[\s_\-0-9]*$', stem): return True
    if re.match(r'^[0-9a-f]{8}-[0-9a-f]{4}-', stem): return True
    if re.match(r'^\d{10,}$', stem): return True
    return False

def _pruefen_passwort(passwort: str) -> bool:
    cfg = _get_config()
    if not cfg.get("passwort_aktiv") or not cfg.get("passwort_hash"):
        return True
    return hashlib.sha256(passwort.encode()).hexdigest() == cfg["passwort_hash"]

# ── Textextraktion ────────────────────────────────────────────

def _extrahiere_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()

    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                seiten = [p.extract_text() or "" for p in pdf.pages[:8]]
                return "\n".join(seiten)[:MAX_TEXT_LEN]
        except ImportError:
            pass
        try:
            import PyPDF2
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages[:8]:
                    text += page.extract_text() or ""
            return text[:MAX_TEXT_LEN]
        except Exception:
            return ""

    if ext in (".docx",):
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])[:MAX_TEXT_LEN]
        except Exception:
            return ""

    if ext in (".xlsx", ".xlsm", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets[:3]:
                for row in ws.iter_rows(max_row=30, values_only=True):
                    rows.append(" | ".join([str(c) for c in row if c is not None]))
            return "\n".join(rows)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    if ext in (".txt", ".csv", ".md", ".rtf"):
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:MAX_TEXT_LEN]
        except Exception:
            return ""

    if ext in (".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif", ".bmp"):
        try:
            import pytesseract
            from PIL import Image, ImageOps
            img = Image.open(filepath)
            img = ImageOps.exif_transpose(img)  # EXIF-Rotation korrigieren (Handy-Fotos!)
            text = pytesseract.image_to_string(img, lang="deu+eng")
            return text[:MAX_TEXT_LEN]
        except ImportError:
            return f"[Scan/Bild: {Path(filepath).name}]"
        except Exception:
            return f"[Scan/Bild: {Path(filepath).name}]"

    if ext in (".odt", ".ods", ".odp"):
        try:
            from odf.opendocument import load
            from odf.text import P
            doc = load(filepath)
            texts = [p.__str__().strip() for p in doc.getElementsByType(P) if p.__str__().strip()]
            return "\n".join(texts)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            return "\n".join(texts)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    return ""


# ── KI-Kategorisierung v6 ─────────────────────────────────────

def _extrahiere_pipe_zeile(antwort: str) -> list | None:
    """
    Ultra-permissiver Scanner: findet jede Zeile mit mindestens 3 Pipes (|).
    Entfernt vorher Markdown-Zeichen und ERGEBNIS-Prefix.
    Kein erzwungener Jahr-Check – wird danach per regex geprüft.
    """
    for zeile in reversed(antwort.splitlines()):
        zeile = re.sub(r'[*_`>#]', '', zeile).strip()
        zeile = re.sub(r'^(?:ERGEBNIS|ergebnis)\s*:?\s*', '', zeile, flags=re.IGNORECASE)
        t = [x.strip() for x in zeile.split("|")]
        if len(t) >= 4 and all(len(x) > 0 for x in t[:4]):
            return t
    return None


def _ki_kategorisiere(filename: str, text: str, provider) -> dict:
    """
    Dokumenten-Kategorisierung v6:
    - Kein CoT bei leerem Text (kurzer Direkt-Prompt)
    - CoT bei vorhandenem Text mit Anti-Halluzinations-Regeln
    - Komplette KI-Antwort im Terminal sichtbar (kein Abschneiden)
    - _extrahiere_pipe_zeile: findet Ergebnis unabhängig von Markdown/Leerzeichen
    """
    ist_kryptisch = _is_kryptisch(filename)
    endung        = Path(filename).suffix.lower()
    jahr_aktuell  = str(datetime.now().year)
    ts_fallback   = datetime.now().strftime('%H%M%S')
    kein_text     = not text or len(text.strip()) < 20

    if kein_text:
        prompt = f"""Kategorisiere diese Datei. Kein Text vorhanden.
Antworte NUR mit dieser einen Zeile (nichts sonst):
Unsortiert|Allgemein|{jahr_aktuell}|Unbekanntes_Dokument_{ts_fallback}{endung}"""
    else:
        prompt = f"""Du bist ein professioneller Dokumenten-Analyst.

Originalname: {filename}
OCR-Text:
---
{text[:2500]}
---

ANALYSIERE IN 3 SCHRITTEN:
SCHRITT 1 – EXTRAKTION: Welche Firmen stehen im Briefkopf? Welches Datum ist erkennbar?
SCHRITT 2 – VALIDIERUNG: Wer ist der RECHTLICHE ABSENDER (oben links)? Firmen die nur als Partner erwähnt werden sind NICHT der Absender. Erfinde keine Absender!
SCHRITT 3 – ENTSCHEIDUNG: Wähle Hauptkategorie, Unterkategorie (Absender) und einen logischen Dateinamen.

Hauptkategorie NUR aus: / Rechnungen / Verträge / Versicherung / Steuern / Behörden / Medizin / Privat / Finanzen/ Arbeit / Immobilien / Fahrzeuge / Bildung / PKW / Unsortiert / 

Schreibe als LETZTE ZEILE deiner Antwort (und nichts danach):
ERGEBNIS: Hauptkategorie|Unterkategorie|Jahr|Absender_Typ_Datum{endung}"""

    try:
        antwort_roh = provider.chat([{"role": "user", "content": prompt}]).strip()

        # DEBUG: Komplette KI-Antwort anzeigen (kein Abschneiden!)
        print(f"\n{'='*60}\nDEBUG '{filename}':\n{antwort_roh}\n{'='*60}\n")

        antwort = re.sub(r'```[^\n]*\n?', '', antwort_roh).strip()

        teile = _extrahiere_pipe_zeile(antwort)
        if not teile:
            raise ValueError("Keine Zeile mit 3+ Trennzeichen '|' gefunden.")

        kat       = _sanitize(teile[0])
        sub       = _sanitize(teile[1])
        jahr      = teile[2].strip()
        dateiname = _sanitize_filename(teile[3].strip(), endung)

        print(f"DEBUG: ✅ KI Format erkannt: {kat}/{sub}/{jahr}/{dateiname}")

        return {
            "kategorie":      kat or "Unsortiert",
            "unterkategorie": sub or "Allgemein",
            "jahr":           jahr if re.match(r'^20\d{2}$', jahr) else jahr_aktuell,
            "dateiname":      dateiname,
            "umbenannt":      ist_kryptisch or True,
        }

    except Exception as e:
        print(f"DEBUG: ❌ Fehler bei '{filename}': {e}")

    # Fallback
    ts       = datetime.now().strftime("%Y%m%d")
    stem     = Path(filename).stem[:40]
    fallback = f"{ts}_{stem}{endung}" if not ist_kryptisch else f"Dokument_{ts}{endung}"
    print(f"DEBUG: ⚠️ Fallback-Name: {fallback}")

    return {
        "kategorie":      "Unsortiert",
        "unterkategorie": "Allgemein",
        "jahr":           jahr_aktuell,
        "dateiname":      fallback,
        "umbenannt":      False,
    }


# ── Öffentliche Skill-Funktionen ──────────────────────────────

def dms_import_scan() -> str:
    """
    Zeigt alle Dateien im Import-Ordner die auf Einsortierung warten.
    Listet Dateinamen und Dateitypen auf ohne KI zu verwenden.
    Beispiel: dms_import_scan()
    """
    _init_dirs()
    import_dir = _get_import_dir()
    dateien = [
        f for f in os.listdir(import_dir)
        if os.path.isfile(os.path.join(import_dir, f))
        and Path(f).suffix.lower() in SUPPORTED_EXTS
    ]
    if not dateien:
        return "📂 Import-Ordner ist leer."
    liste = "\n".join([f"  • {f}" for f in sorted(dateien)])
    return f"📥 {len(dateien)} Dokument(e) im Import-Ordner:\n{liste}"


def dms_einsortieren(provider=None) -> str:
    """KI liest, benennt und archiviert alle Dokumente im Import-Ordner."""
    _init_dirs()
    import_dir = _get_import_dir()
    archiv_dir = _get_archiv_dir()

    if provider is None:
        try:
            from providers import select_provider
            _, provider = select_provider("auto")
        except Exception as e:
            return f"❌ Kein KI-Provider verfügbar: {e}"

    dateien = [
        f for f in os.listdir(import_dir)
        if os.path.isfile(os.path.join(import_dir, f))
        and Path(f).suffix.lower() in SUPPORTED_EXTS
    ]

    if not dateien:
        return "📂 Import-Ordner ist leer."

    meta    = _load_meta()
    bericht = []

    for dateiname in sorted(dateien):
        quell_pfad = os.path.join(import_dir, dateiname)
        endung     = Path(dateiname).suffix.lower()
        datei_hash = _berechne_hash(quell_pfad)

        # Duplikat-Check
        if datei_hash and datei_hash in meta.get("hashes", {}):
            vorhandener = meta["hashes"][datei_hash]
            os.remove(quell_pfad)
            bericht.append(f"♻️ Duplikat: **{dateiname}** → identisch mit {vorhandener}")
            continue

        text = _extrahiere_text(quell_pfad)
        clean_text = text.replace('\n', ' ')
        print(f"DEBUG: OCR-Text für {dateiname}: {clean_text[:120]}...")

        ki_info    = _ki_kategorisiere(dateiname, text, provider)
        neuer_name = ki_info["dateiname"]
        if not neuer_name.lower().endswith(endung):
            neuer_name = Path(neuer_name).stem + endung

        ziel_ordner = os.path.join(
            archiv_dir,
            ki_info["kategorie"],
            ki_info["unterkategorie"],
            ki_info["jahr"]
        )

        if not os.path.abspath(ziel_ordner).startswith(os.path.abspath(archiv_dir)):
            bericht.append(f"❌ Sicherheitswarnung: {dateiname} – ungültiger Zielpfad.")
            continue

        os.makedirs(ziel_ordner, exist_ok=True)
        wunsch_ziel  = os.path.join(ziel_ordner, neuer_name)
        finaler_pfad = _naechste_version(wunsch_ziel)
        finaler_name = os.path.basename(finaler_pfad)

        shutil.move(quell_pfad, finaler_pfad)

        rel_pfad = os.path.relpath(finaler_pfad, archiv_dir).replace("\\", "/")

        if "hashes"    not in meta: meta["hashes"]    = {}
        if "dokumente" not in meta: meta["dokumente"]  = {}

        if datei_hash:
            meta["hashes"][datei_hash] = rel_pfad
        meta["dokumente"][rel_pfad] = {
            "original":   dateiname,
            "kategorie":  ki_info["kategorie"],
            "sub":        ki_info["unterkategorie"],
            "jahr":       ki_info["jahr"],
            "hash":       datei_hash,
            "groesse":    os.path.getsize(finaler_pfad),
            "archiviert": datetime.now().isoformat(),
            "umbenannt":  ki_info.get("umbenannt", False),
        }

        umbenennt = f" ✏️ ← '{dateiname}'" if ki_info.get("umbenannt") and finaler_name != dateiname else ""
        version   = " 🆙 (neue Version)" if finaler_pfad != wunsch_ziel else ""
        bericht.append(f"✅ **{finaler_name}**{umbenennt}{version}\n   → {ki_info['kategorie']} / {ki_info['unterkategorie']} / {ki_info['jahr']}")

    _save_meta(meta)
    return f"📁 {len(bericht)} Dokument(e) verarbeitet:\n\n" + "\n\n".join(bericht)


def dms_loeschen(pfad_relativ: str, passwort: str = "") -> str:
    """Löscht eine archivierte Datei und entfernt sie aus meta.json."""
    _init_dirs()

    if not _pruefen_passwort(passwort):
        return "❌ Falsches Passwort."

    archiv_dir = _get_archiv_dir()
    voll_pfad  = os.path.abspath(os.path.join(archiv_dir, pfad_relativ))

    if not voll_pfad.startswith(os.path.abspath(archiv_dir)):
        return "❌ Ungültiger Pfad."

    meta = _load_meta()
    norm = pfad_relativ.replace("\\", "/")

    if not os.path.isfile(voll_pfad):
        _entferne_aus_meta(meta, norm)
        _save_meta(meta)
        return f"ℹ️ Datei nicht gefunden, aber aus Index entfernt."

    os.remove(voll_pfad)

    try:
        eltern = Path(voll_pfad).parent
        while str(eltern) != str(archiv_dir):
            if not any(eltern.iterdir()):
                eltern.rmdir()
                eltern = eltern.parent
            else:
                break
    except Exception:
        pass

    _entferne_aus_meta(meta, norm)
    _save_meta(meta)
    return f"🗑️ Gelöscht: {pfad_relativ}"


def dms_verschieben(pfad_relativ: str, neue_kategorie: str, neue_unterkategorie: str = "", passwort: str = "") -> dict:
    """
    Verschiebt eine archivierte Datei in eine andere Kategorie/Unterkategorie.
    Behält den Dateinamen und das Jahr bei.
    Gibt {"ok": True, "neuer_pfad": "..."} oder {"ok": False, "error": "..."} zurück.
    """
    _init_dirs()

    if not _pruefen_passwort(passwort):
        return {"ok": False, "error": "Falsches Passwort."}

    archiv_dir = _get_archiv_dir()
    voll_pfad  = os.path.abspath(os.path.join(archiv_dir, pfad_relativ))

    if not voll_pfad.startswith(os.path.abspath(archiv_dir)):
        return {"ok": False, "error": "Ungültiger Quellpfad."}

    if not os.path.isfile(voll_pfad):
        return {"ok": False, "error": "Datei nicht gefunden."}

    # Kategorie säubern
    neue_kategorie    = _sanitize(neue_kategorie)
    neue_unterkategorie = _sanitize(neue_unterkategorie) if neue_unterkategorie else ""

    # Jahr und Dateiname aus dem alten Pfad übernehmen
    teile     = pfad_relativ.replace("\\", "/").split("/")
    dateiname = teile[-1]
    altes_jahr = teile[2] if len(teile) >= 4 else str(datetime.now().year)
    alte_sub   = teile[1] if len(teile) >= 3 else "Allgemein"

    # Unterkategorie: falls nicht angegeben, alte behalten
    if not neue_unterkategorie:
        neue_unterkategorie = alte_sub

    ziel_ordner = os.path.join(archiv_dir, neue_kategorie, neue_unterkategorie, altes_jahr)
    if not os.path.abspath(ziel_ordner).startswith(os.path.abspath(archiv_dir)):
        return {"ok": False, "error": "Ungültiger Zielpfad."}

    os.makedirs(ziel_ordner, exist_ok=True)
    ziel_pfad   = _naechste_version(os.path.join(ziel_ordner, dateiname))
    shutil.move(voll_pfad, ziel_pfad)

    # Leere Quell-Ordner aufräumen
    try:
        eltern = Path(voll_pfad).parent
        while str(eltern) != str(archiv_dir):
            if not any(eltern.iterdir()):
                eltern.rmdir()
                eltern = eltern.parent
            else:
                break
    except Exception:
        pass

    # Meta aktualisieren
    meta     = _load_meta()
    alter_rel = pfad_relativ.replace("\\", "/")
    neuer_rel = os.path.relpath(ziel_pfad, archiv_dir).replace("\\", "/")

    # Eintrag umbenennen
    dok = meta.get("dokumente", {}).pop(alter_rel, {})
    if dok:
        dok["kategorie"]  = neue_kategorie
        dok["sub"]        = neue_unterkategorie
        meta.setdefault("dokumente", {})[neuer_rel] = dok

    # Hash-Zeiger aktualisieren
    for h, p in list(meta.get("hashes", {}).items()):
        if p.replace("\\", "/") == alter_rel:
            meta["hashes"][h] = neuer_rel

    _save_meta(meta)
    return {"ok": True, "neuer_pfad": neuer_rel, "dateiname": os.path.basename(ziel_pfad)}


def _entferne_aus_meta(meta: dict, rel_pfad: str):
    norm = rel_pfad.replace("\\", "/")
    for k in list(meta.get("dokumente", {}).keys()):
        if k.replace("\\", "/") == norm:
            del meta["dokumente"][k]
    for h in list(meta.get("hashes", {}).keys()):
        if meta["hashes"][h].replace("\\", "/") == norm:
            del meta["hashes"][h]


def dms_suchen(suchbegriff: str) -> str:
    """
    Durchsucht das DMS-Archiv nach Dateien die den Suchbegriff im Namen enthalten.
    Gibt Dateipfad und Größe aller Treffer zurück.
    Beispiel: dms_suchen(suchbegriff="Rechnung 2024")
    """
    _init_dirs()
    archiv_dir = _get_archiv_dir()
    treffer    = []
    q          = suchbegriff.lower()
    for root, _, files in os.walk(archiv_dir):
        for f in files:
            rel = os.path.relpath(os.path.join(root, f), archiv_dir)
            if q in rel.lower():
                sz   = os.path.getsize(os.path.join(root, f))
                sz_s = f"{sz//1024} KB" if sz > 1024 else f"{sz} B"
                treffer.append(f"  📄 {rel}  ({sz_s})")
    if not treffer:
        return f"🔍 Keine Treffer für: '{suchbegriff}'"
    return f"🔍 {len(treffer)} Treffer:\n" + "\n".join(sorted(treffer))


def dms_archiv_uebersicht() -> str:
    """
    Zeigt eine kompakte Übersicht aller Kategorien im DMS-Archiv mit Dokumentenanzahl.
    Beispiel: dms_archiv_uebersicht()
    """
    _init_dirs()
    archiv_dir = _get_archiv_dir()
    struktur   = {}
    for root, _, files in os.walk(archiv_dir):
        for f in files:
            if Path(f).suffix.lower() in SUPPORTED_EXTS:
                rel  = os.path.relpath(root, archiv_dir)
                kat  = rel.split(os.sep)[0] if rel != "." else "Unsortiert"
                struktur[kat] = struktur.get(kat, 0) + 1
    if not struktur:
        return "📁 Archiv leer."
    gesamt = sum(struktur.values())
    zeilen = [f"📁 {gesamt} Dokumente gesamt:"]
    for kat, n in sorted(struktur.items()):
        zeilen.append(f"  📂 {kat}: {n}")
    return "\n".join(zeilen)


def dms_stats() -> dict:
    """
    Gibt detaillierte Statistiken über das DMS-Archiv zurück: Gesamtanzahl, Größe in MB, Kategorien, Import-Warteschlange.
    Beispiel: dms_stats()
    """
    _init_dirs()
    archiv_dir   = _get_archiv_dir()
    import_dir   = _get_import_dir()
    gesamt       = 0
    groesse      = 0
    kategorien   = {}

    for root, _, files in os.walk(archiv_dir):
        for f in files:
            if Path(f).suffix.lower() in SUPPORTED_EXTS:
                gesamt  += 1
                groesse += os.path.getsize(os.path.join(root, f))
                rel      = os.path.relpath(root, archiv_dir)
                teile    = rel.split(os.sep)
                kat      = teile[0] if teile and teile[0] != "." else "Unsortiert"
                kategorien[kat] = kategorien.get(kat, 0) + 1

    import_count = len([
        f for f in os.listdir(import_dir)
        if os.path.isfile(os.path.join(import_dir, f))
    ]) if os.path.exists(import_dir) else 0

    cfg = _get_config()
    return {
        "gesamt":         gesamt,
        "groesse_mb":     round(groesse / (1024 * 1024), 2),
        "kategorien":     kategorien,
        "import_count":   import_count,
        "archiv_pfad":    cfg["archiv_pfad"],
        "import_pfad":    cfg["import_pfad"],
        "passwort_aktiv": cfg.get("passwort_aktiv", False),
    }


def dms_archiv_baum() -> list:
    """
    Gibt die vollständige Ordnerstruktur des DMS-Archivs als verschachtelte Liste zurück.
    Nützlich um alle Kategorien und Unterkategorien auf einen Blick zu sehen.
    Beispiel: dms_archiv_baum()
    """
    _init_dirs()
    archiv_dir = _get_archiv_dir()
    baum       = {}

    for root, _, files in os.walk(archiv_dir):
        for f in files:
            if Path(f).suffix.lower() in SUPPORTED_EXTS:
                voll    = os.path.join(root, f)
                rel     = os.path.relpath(voll, archiv_dir)
                teile   = rel.split(os.sep)
                kat     = teile[0] if len(teile) > 0 else "Unsortiert"
                sub     = teile[1] if len(teile) > 1 else ""
                jahr    = teile[2] if len(teile) > 2 else ""
                groesse = os.path.getsize(voll)
                mtime   = datetime.fromtimestamp(os.path.getmtime(voll)).strftime("%d.%m.%Y")

                if kat not in baum:
                    baum[kat] = {}
                key = f"{sub}/{jahr}" if sub else "Allgemein"
                if key not in baum[kat]:
                    baum[kat][key] = []
                baum[kat][key].append({
                    "name":    f,
                    "pfad":    rel.replace("\\", "/"),
                    "groesse": groesse,
                    "datum":   mtime,
                    "ext":     Path(f).suffix.lower().lstrip("."),
                })

    ergebnis = []
    for kat, subs in sorted(baum.items()):
        gesamt = sum(len(d) for d in subs.values())
        obj    = {"name": kat, "count": gesamt, "subs": []}
        for sub, dateien in sorted(subs.items()):
            obj["subs"].append({
                "name":    sub,
                "dateien": sorted(dateien, key=lambda x: x["name"])
            })
        ergebnis.append(obj)
    return ergebnis


def dms_pfad_setzen(archiv_pfad: str, import_pfad: str = "", passwort: str = "", passwort_neu: str = "") -> str:
    """
    Legt die Archiv- und Import-Pfade für das DMS fest. Einmalig ausführen beim ersten Start.
    Optional: Passwortschutz für sensible Operationen aktivieren.
    Beispiel: dms_pfad_setzen(archiv_pfad="/home/user/archiv", import_pfad="/home/user/import")
    """
    cfg = _get_config()
    if cfg.get("passwort_aktiv") and cfg.get("passwort_hash"):
        if not _pruefen_passwort(passwort):
            return "❌ Falsches Passwort."

    if archiv_pfad:
        p = os.path.expandvars(os.path.expanduser(archiv_pfad))
        try:
            os.makedirs(p, exist_ok=True)
            cfg["archiv_pfad"] = os.path.abspath(p)
        except Exception as e:
            return f"❌ Archiv-Pfad ungültig: {e}"

    if import_pfad:
        p = os.path.expandvars(os.path.expanduser(import_pfad))
        try:
            os.makedirs(p, exist_ok=True)
            cfg["import_pfad"] = os.path.abspath(p)
        except Exception as e:
            return f"❌ Import-Pfad ungültig: {e}"

    if passwort_neu:
        cfg["passwort_hash"]  = hashlib.sha256(passwort_neu.encode()).hexdigest()
        cfg["passwort_aktiv"] = True

    _save_config(cfg)
    return f"✅ Gespeichert:\n  Archiv: {cfg['archiv_pfad']}\n  Import: {cfg['import_pfad']}"


def dms_passwort_entfernen(passwort: str) -> str:
    """
    Entfernt den Passwortschutz vom DMS-Archiv. Aktuelles Passwort zur Bestätigung erforderlich.
    Beispiel: dms_passwort_entfernen(passwort="meinPasswort")
    """
    cfg = _get_config()
    if cfg.get("passwort_aktiv"):
        if not _pruefen_passwort(passwort):
            return "❌ Falsches Passwort."
    cfg["passwort_aktiv"] = False
    cfg["passwort_hash"]  = ""
    _save_config(cfg)
    return "✅ Passwortschutz entfernt."


AVAILABLE_SKILLS = [
    dms_import_scan,
    dms_einsortieren,
    dms_suchen,
    dms_archiv_uebersicht,
    dms_loeschen,
    dms_stats,
]
