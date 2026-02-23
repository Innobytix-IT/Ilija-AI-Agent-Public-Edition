"""
DMS – Dokumentenmanagementsystem für Ilija Public Edition
==========================================================
Unterstützte Formate: PDF, DOCX, DOC, XLSX, XLS, TXT, CSV,
                      JPG, PNG, WEBP, TIFF, BMP, HEIC,
                      MD, RTF, ODP, ODS, ODT
KI-gestützte automatische Kategorisierung & Archivierung
Duplikat-Erkennung via SHA-256, automatische Versionierung
"""

import os
import re
import shutil
import hashlib
import json
from pathlib import Path
from datetime import datetime

# ── Konfiguration ────────────────────────────────────────────────
DMS_BASE     = os.path.abspath("data/dms")
IMPORT_DIR   = os.path.join(DMS_BASE, "import")
ARCHIV_DIR   = os.path.join(DMS_BASE, "archiv")
META_FILE    = os.path.join(DMS_BASE, "meta.json")

# Maximale Zeichen die an die KI gesendet werden
MAX_TEXT_LEN = 3000

# Unterstützte Formate
SUPPORTED_EXTS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".xlsm",
    ".txt", ".csv", ".md", ".rtf",
    ".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif",
    ".bmp", ".heic", ".heif",
    ".odt", ".ods", ".odp", ".pptx", ".ppt",
}

# ── Hilfsfunktionen ──────────────────────────────────────────────

def _init_dirs():
    """Erstellt alle nötigen Verzeichnisse."""
    os.makedirs(IMPORT_DIR, exist_ok=True)
    os.makedirs(ARCHIV_DIR, exist_ok=True)
    if not os.path.exists(META_FILE):
        _save_meta({})

def _load_meta() -> dict:
    try:
        with open(META_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def _save_meta(meta: dict):
    with open(META_FILE, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

def _berechne_hash(filepath: str) -> str:
    """SHA-256 Fingerabdruck einer Datei."""
    hasher = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
    except Exception:
        return ""
    return hasher.hexdigest()

def _naechste_version(ziel_pfad: str) -> str:
    """Gibt freien Versionspfad zurück: datei_v2.pdf, datei_v3.pdf ..."""
    if not os.path.exists(ziel_pfad):
        return ziel_pfad
    basis, endung = os.path.splitext(ziel_pfad)
    # Eventuell schon vorhandenes _vN entfernen
    basis = re.sub(r'_v\d+$', '', basis)
    counter = 2
    while True:
        kandidat = f"{basis}_v{counter}{endung}"
        if not os.path.exists(kandidat):
            return kandidat
        counter += 1

def _sanitize(text: str) -> str:
    """Bereinigt Zeichenketten für Dateinamen."""
    text = text.strip()
    text = re.sub(r'[\\/*?:"<>|]', '', text)
    text = re.sub(r'\s+', '_', text)
    return text[:80]

# ── Textextraktion ───────────────────────────────────────────────

def _extrahiere_text(filepath: str) -> str:
    """Extrahiert lesbaren Text aus allen unterstützten Formaten."""
    ext = Path(filepath).suffix.lower()

    # ── PDF ──────────────────────────────────────────────────────
    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                seiten = [p.extract_text() or "" for p in pdf.pages[:8]]
                return "\n".join(seiten)[:MAX_TEXT_LEN]
        except ImportError:
            pass
        try:
            import PyPDF2
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages[:8]:
                    text += page.extract_text() or ""
            return text[:MAX_TEXT_LEN]
        except Exception:
            return ""

    # ── Word ─────────────────────────────────────────────────────
    if ext in (".docx",):
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])[:MAX_TEXT_LEN]
        except Exception:
            return ""

    # ── Excel ────────────────────────────────────────────────────
    if ext in (".xlsx", ".xlsm", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets[:3]:
                for row in ws.iter_rows(max_row=30, values_only=True):
                    rows.append(" | ".join([str(c) for c in row if c is not None]))
            return "\n".join(rows)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    # ── Text / CSV / Markdown ────────────────────────────────────
    if ext in (".txt", ".csv", ".md", ".rtf"):
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:MAX_TEXT_LEN]
        except Exception:
            return ""

    # ── Bilder (OCR) ─────────────────────────────────────────────
    if ext in (".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif", ".bmp"):
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img, lang="deu+eng")
            return text[:MAX_TEXT_LEN]
        except ImportError:
            return f"[Bild: {Path(filepath).name} – pytesseract nicht installiert]"
        except Exception:
            return f"[Bild: {Path(filepath).name} – OCR fehlgeschlagen]"

    # ── OpenDocument ─────────────────────────────────────────────
    if ext in (".odt", ".ods", ".odp"):
        try:
            from odf.opendocument import load
            from odf.text import P
            doc = load(filepath)
            texts = []
            for p in doc.getElementsByType(P):
                t = p.__str__()
                if t.strip():
                    texts.append(t.strip())
            return "\n".join(texts)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    # ── PowerPoint ───────────────────────────────────────────────
    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)[:MAX_TEXT_LEN]
        except Exception:
            return ""

    return ""


# ── KI-Kategorisierung ───────────────────────────────────────────

def _ki_kategorisiere(filename: str, text: str, provider) -> dict:
    """
    Lässt die KI Kategorie, Jahr, Unterkategorie und Dateinamen bestimmen.
    Rückgabe: {"kategorie": ..., "unterkategorie": ..., "jahr": ..., "dateiname": ...}
    """
    prompt = f"""Du bist ein intelligentes Dokumentenarchiv-System.
Analysiere das folgende Dokument und kategorisiere es präzise.

Originaldateiname: {filename}

Dokumenteninhalt (Auszug):
---
{text if text else '[Inhalt konnte nicht extrahiert werden – nutze den Dateinamen]'}
---

Antworte NUR in diesem exakten Format (eine Zeile, Trennzeichen |):
KATEGORIE|UNTERKATEGORIE|JAHR|NEUER_DATEINAME

Regeln:
- KATEGORIE: Hauptkategorie (z.B. Rechnungen, Verträge, Versicherung, Behörden, Steuern, Medizin, Privat, Finanzen, Arbeit, Immobilien, Fahrzeuge, Bildung)
- UNTERKATEGORIE: Genauere Einordnung (z.B. Telekom, HUK-Coburg, Finanzamt, Lohnabrechnung)
- JAHR: 4-stellige Jahreszahl aus dem Dokument, oder aktuelles Jahr wenn unklar
- NEUER_DATEINAME: Beschreibender Dateiname OHNE Leerzeichen, MIT Dateiendung (z.B. Rechnung_Strom_Maerz.pdf)

Nur eine Zeile, keine Erklärungen, kein Markdown."""

    try:
        antwort = provider.chat([{"role": "user", "content": prompt}]).strip()
        # Nur erste Zeile nehmen
        antwort = antwort.split("\n")[0].strip()
        teile = antwort.split("|")
        if len(teile) >= 4:
            return {
                "kategorie":     _sanitize(teile[0]),
                "unterkategorie": _sanitize(teile[1]),
                "jahr":          _sanitize(teile[2]) if teile[2].isdigit() else str(datetime.now().year),
                "dateiname":     _sanitize(teile[3]),
            }
    except Exception:
        pass

    # Fallback: Originaldateiname behalten
    return {
        "kategorie":      "Unsortiert",
        "unterkategorie": "Allgemein",
        "jahr":           str(datetime.now().year),
        "dateiname":      Path(filename).name,
    }


# ── Öffentliche Skill-Funktionen ─────────────────────────────────

def dms_import_scan() -> str:
    """
    Zeigt alle Dateien im Import-Ordner die noch nicht archiviert wurden.
    Bereit zum Einsortieren.
    """
    _init_dirs()
    dateien = [
        f for f in os.listdir(IMPORT_DIR)
        if os.path.isfile(os.path.join(IMPORT_DIR, f))
        and Path(f).suffix.lower() in SUPPORTED_EXTS
    ]
    if not dateien:
        return "📂 Import-Ordner ist leer. Keine neuen Dokumente zum Einsortieren."
    
    liste = "\n".join([f"  • {f}" for f in sorted(dateien)])
    return f"📥 {len(dateien)} Dokument(e) im Import-Ordner:\n{liste}\n\nSage 'Dokumente einsortieren' um sie zu archivieren."


def dms_einsortieren(provider=None) -> str:
    """
    Analysiert alle Dateien im Import-Ordner per KI und archiviert sie.
    Duplikate werden erkannt, neue Versionen automatisch erstellt.
    """
    _init_dirs()

    if provider is None:
        try:
            from providers import select_provider
            _, provider = select_provider("auto")
        except Exception as e:
            return f"❌ Kein KI-Provider verfügbar: {e}"

    dateien = [
        f for f in os.listdir(IMPORT_DIR)
        if os.path.isfile(os.path.join(IMPORT_DIR, f))
        and Path(f).suffix.lower() in SUPPORTED_EXTS
    ]

    if not dateien:
        return "📂 Import-Ordner ist leer. Keine Dokumente zum Einsortieren."

    meta = _load_meta()
    bericht = []

    for dateiname in sorted(dateien):
        quell_pfad = os.path.join(IMPORT_DIR, dateiname)
        endung     = Path(dateiname).suffix.lower()

        # Hash berechnen
        datei_hash = _berechne_hash(quell_pfad)

        # Duplikat-Check über alle bereits archivierten Dateien
        if datei_hash and datei_hash in meta.get("hashes", {}):
            vorhandener = meta["hashes"][datei_hash]
            os.remove(quell_pfad)
            bericht.append(f"♻️ Duplikat gelöscht: **{dateiname}** (identisch mit {vorhandener})")
            continue

        # Text extrahieren
        text = _extrahiere_text(quell_pfad)

        # KI-Kategorisierung
        ki_info = _ki_kategorisiere(dateiname, text, provider)

        # Sicherstellen dass Dateiendung erhalten bleibt
        neuer_name = ki_info["dateiname"]
        if not neuer_name.lower().endswith(endung):
            neuer_name = Path(neuer_name).stem + endung

        # Zielpfad aufbauen
        ziel_ordner = os.path.join(
            ARCHIV_DIR,
            ki_info["kategorie"],
            ki_info["unterkategorie"],
            ki_info["jahr"]
        )

        # Sicherheitscheck: Path Traversal verhindern
        if not os.path.abspath(ziel_ordner).startswith(os.path.abspath(ARCHIV_DIR)):
            bericht.append(f"❌ Sicherheitswarnung bei {dateiname}: Ungültiger Zielpfad.")
            continue

        os.makedirs(ziel_ordner, exist_ok=True)
        wunsch_ziel = os.path.join(ziel_ordner, neuer_name)

        # Versionierung falls Name bereits existiert
        finaler_pfad = _naechste_version(wunsch_ziel)
        finaler_name = os.path.basename(finaler_pfad)

        # Verschieben
        shutil.move(quell_pfad, finaler_pfad)

        # Meta aktualisieren
        rel_pfad = os.path.relpath(finaler_pfad, ARCHIV_DIR)
        if "hashes" not in meta:
            meta["hashes"] = {}
        if "dokumente" not in meta:
            meta["dokumente"] = {}

        meta["hashes"][datei_hash] = rel_pfad
        meta["dokumente"][rel_pfad] = {
            "original":   dateiname,
            "kategorie":  ki_info["kategorie"],
            "sub":        ki_info["unterkategorie"],
            "jahr":       ki_info["jahr"],
            "hash":       datei_hash,
            "groesse":    os.path.getsize(finaler_pfad),
            "archiviert": datetime.now().isoformat(),
        }

        if finaler_pfad != wunsch_ziel:
            version = re.search(r'_v(\d+)', finaler_name)
            v_str   = f" (Version {version.group(1)})" if version else ""
            bericht.append(f"🆙 Neue Version{v_str}: **{dateiname}** → {rel_pfad}")
        else:
            bericht.append(f"✅ Archiviert: **{dateiname}** → {rel_pfad}")

    _save_meta(meta)

    zusammenfassung = f"📁 Ergebnis: {len(bericht)} Dokument(e) verarbeitet.\n\n"
    return zusammenfassung + "\n".join(bericht)


def dms_suchen(suchbegriff: str) -> str:
    """
    Durchsucht das Archiv nach Dateien deren Name oder Pfad den Suchbegriff enthält.
    """
    _init_dirs()
    treffer = []
    suchbegriff_lower = suchbegriff.lower()

    for root, _, files in os.walk(ARCHIV_DIR):
        for f in files:
            rel = os.path.relpath(os.path.join(root, f), ARCHIV_DIR)
            if suchbegriff_lower in rel.lower():
                groesse = os.path.getsize(os.path.join(root, f))
                groesse_str = f"{groesse // 1024} KB" if groesse > 1024 else f"{groesse} B"
                treffer.append(f"  📄 {rel}  ({groesse_str})")

    if not treffer:
        return f"🔍 Keine Dokumente gefunden für: '{suchbegriff}'"

    return f"🔍 {len(treffer)} Treffer für '{suchbegriff}':\n" + "\n".join(sorted(treffer))


def dms_archiv_uebersicht() -> str:
    """
    Zeigt die gesamte Archivstruktur als Baum mit Dateianzahl pro Kategorie.
    """
    _init_dirs()
    struktur = {}

    for root, _, files in os.walk(ARCHIV_DIR):
        for f in files:
            if Path(f).suffix.lower() in SUPPORTED_EXTS:
                rel = os.path.relpath(root, ARCHIV_DIR)
                teile = rel.split(os.sep)
                kat = teile[0] if teile else "Unsortiert"
                struktur[kat] = struktur.get(kat, 0) + 1

    if not struktur:
        return "📁 Das Archiv ist leer."

    gesamt = sum(struktur.values())
    zeilen = [f"📁 Archiv-Übersicht ({gesamt} Dokumente gesamt):"]
    for kat, anzahl in sorted(struktur.items()):
        zeilen.append(f"  📂 {kat}  –  {anzahl} Dokument(e)")

    return "\n".join(zeilen)


def dms_stats() -> dict:
    """Gibt Statistiken für die Web GUI zurück."""
    _init_dirs()
    meta    = _load_meta()
    docs    = meta.get("dokumente", {})
    gesamt  = len(docs)
    groesse = sum(v.get("groesse", 0) for v in docs.values())

    kategorien = {}
    for v in docs.values():
        k = v.get("kategorie", "Unsortiert")
        kategorien[k] = kategorien.get(k, 0) + 1

    import_count = len([
        f for f in os.listdir(IMPORT_DIR)
        if os.path.isfile(os.path.join(IMPORT_DIR, f))
    ]) if os.path.exists(IMPORT_DIR) else 0

    return {
        "gesamt":       gesamt,
        "groesse_mb":   round(groesse / (1024 * 1024), 2),
        "kategorien":   kategorien,
        "import_count": import_count,
    }


def dms_archiv_baum() -> list:
    """Gibt den Archivbaum als verschachtelte Liste für die Web GUI zurück."""
    _init_dirs()
    baum = {}

    for root, _, files in os.walk(ARCHIV_DIR):
        for f in files:
            if Path(f).suffix.lower() in SUPPORTED_EXTS:
                voll    = os.path.join(root, f)
                rel     = os.path.relpath(voll, ARCHIV_DIR)
                teile   = rel.split(os.sep)
                kat     = teile[0] if len(teile) > 0 else "Unsortiert"
                sub     = teile[1] if len(teile) > 1 else ""
                jahr    = teile[2] if len(teile) > 2 else ""
                name    = teile[-1]
                groesse = os.path.getsize(voll)
                mtime   = datetime.fromtimestamp(os.path.getmtime(voll)).strftime("%d.%m.%Y")

                if kat not in baum:
                    baum[kat] = {}
                key = f"{sub}/{jahr}" if sub else "Allgemein"
                if key not in baum[kat]:
                    baum[kat][key] = []
                baum[kat][key].append({
                    "name":     name,
                    "pfad":     rel.replace("\\", "/"),
                    "groesse":  groesse,
                    "datum":    mtime,
                    "ext":      Path(f).suffix.lower().lstrip("."),
                })

    # In Liste umwandeln
    ergebnis = []
    for kat, subs in sorted(baum.items()):
        dateien_gesamt = sum(len(d) for d in subs.values())
        kat_obj = {"name": kat, "count": dateien_gesamt, "subs": []}
        for sub, dateien in sorted(subs.items()):
            kat_obj["subs"].append({"name": sub, "dateien": sorted(dateien, key=lambda x: x["name"])})
        ergebnis.append(kat_obj)

    return ergebnis


# ── Skill-Registry ────────────────────────────────────────────────
AVAILABLE_SKILLS = [
    dms_import_scan,
    dms_einsortieren,
    dms_suchen,
    dms_archiv_uebersicht,
]
